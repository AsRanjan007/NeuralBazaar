"""

NeuralBazaar - AI-Powered Stock Market Intelligence Dashboard

Copyright (c) 2026 Ashutosh Ranjan. All rights reserved.

Licensed under the MIT License. See LICENSE file for details.

Version: 1.0.0

"""

from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slides_data = [
    ("Stock Market Dashboard Architecture", "Overview of system components:\n- Streamlit UI\n- YFinance data fetch\n- Indicator engine\n- Strategy module\n- LSTM forecast module\n- Backtest and reporting"),
    ("Data Flow Diagram", "Data Flow:\n1. User selects ticker/dates/strategy\n2. App fetches market OHLC from YFinance\n3. Compute technical indicators (RSI, MACD, Bollinger)\n4. Strategy generates signals\n5. Backtest calculates profit/trades\n6. LSTM forecast predicts next 7 trading days\n7. UI shows charts and summary"),
    ("Technical Indicators", "RSI: overbought/oversold strength\nMACD: trend momentum & crossovers\nBollinger Bands: volatility + mean reversion"),
    ("Strategy Logic", "RSI Strategy: buy when RSI <30, sell when >70\nMACD Strategy: buy when MACD line crosses above signal line, sell reverse\nBollinger Strategy: buy near lower band, sell near upper band"),
    ("LSTM Forecast Module", "Uses past close+indicators to train sequence model\nAttention layer -> dense output forecast\nPredict next 7 valid trading days (skip weekends/holidays)"),
    ("Streamlit Page Components", "Sidebar inputs, tabs: Summary, Chart, Indicators, Forecast, Strategy\nMetrics, charts, tables, download"),
    ("Sequence Diagram", "User -> App: run with settings\nApp -> YFinance: fetch price\nApp -> indicator module: compute indicators\nApp -> strategy: signal generation\nApp -> backtest: trade simulation\nApp -> LSTM: forecast compute\nApp -> Streamlit: render output"),
    ("Summary & Next Steps", "How to use and improve:\n- add more indices and exchange support\n- include risk metrics (Sharpe, drawdown)\n- persistent data caching\n- model versioning and hyperparameter tuning")
]

for title, body in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    textbox = slide.shapes.placeholders[1]
    textbox.text = body

prs.save('stock_market_dashboard_presentation.pptx')
print('Presentation created:', 'stock_market_dashboard_presentation.pptx')
